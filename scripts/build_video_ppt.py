"""Build a concise blank-template PPTX of generated plant videos, annotated with
model / input-source / motion-type. Embeds the actual mp4s with poster frames."""
import glob, os
import numpy as np
import imageio.v2 as imageio
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

GEN = 'outputs/rerun_2026-07/gen'
TMP = '/tmp/ppt_posters'; os.makedirs(TMP, exist_ok=True)


def g1(pattern):
    m = glob.glob(pattern)
    return m[0] if m else None


# curated: (video_path, title, model, input_source, motion_type)
DYNAMIC = [
    (f'{GEN}/newplant_sway/newplant4/out/output.mp4', 'newplant4 sway', 'Cosmos-Predict2 2B', 'image (GS-render frame) → video', 'wind sway'),
    (f'{GEN}/newplant_sway/newplant9/out/output.mp4', 'newplant9 sway', 'Cosmos-Predict2 2B', 'image (GS-render frame) → video', 'wind sway'),
    (g1(f'{GEN}/motion20/a_leafy_young_tree*wb_gentle/motion.mp4'), 'leafy young tree', 'Wan 2.2 I2V-A14B', 'image (TRELLIS pose) → video', 'gentle breeze'),
    (g1(f'{GEN}/motion20/a_fig_sapling*wb_gentle/motion.mp4'), 'fig sapling', 'Wan 2.2 I2V-A14B', 'image (TRELLIS pose) → video', 'gentle breeze'),
    (f'{GEN}/wind_sway.mp4', 'FK wind sway', 'Ours: 3DGS + FK physics', '3D Gaussians → render video', 'articulated sway'),
    (f'{GEN}/multiview_sway.mp4', 'multi-view sway', 'Ours: 3DGS + FK physics', '3D Gaussians → render video', 'multi-view consistent'),
]
GROWTH = [
    (f'{GEN}/cosmos_growth/vine_climbing/output.mp4', 'vine climbing', 'Cosmos-Predict2 2B', 'image (seed) → video', 'growth'),
    (f'{GEN}/cosmos_growth/branch_extending/output.mp4', 'branch extending', 'Cosmos-Predict2 2B', 'image (seed) → video', 'growth'),
    (f'{GEN}/cosmos_growth/leaves_unfurling/output.mp4', 'leaves unfurling', 'Cosmos-Predict2 2B', 'image (seed) → video', 'growth'),
    (f'{GEN}/cosmos_growth/stem_lengthening/output.mp4', 'stem lengthening', 'Cosmos-Predict2 2B', 'image (seed) → video', 'growth'),
]
STATIC3D = [
    (g1(f'{GEN}/trellis_plants/a_flowering_geranium*/preview.mp4'), 'flowering geranium', 'TRELLIS', 'text → 3D Gaussians → render', 'novel-view orbit'),
    (g1(f'{GEN}/trellis_plants/a_small_basil*/preview.mp4'), 'basil plant', 'TRELLIS', 'text → 3D Gaussians → render', 'novel-view orbit'),
    (f'{GEN}/multiview_orbit.mp4', 'novel-view orbit', 'Ours: 3DGS render', '3D Gaussians → render', 'camera orbit'),
]

INK = RGBColor(0x1a, 0x1a, 0x1a)
SUB = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x2e, 0x7d, 0x32)
WHITE = RGBColor(0xff, 0xff, 0xff)


def poster(path):
    v = imageio.mimread(path, memtest=False)
    f = v[len(v) // 2]
    out = f'{TMP}/{abs(hash(path))}.jpg'
    imageio.imwrite(out, f, quality=85)
    h, w = f.shape[:2]
    return out, w / h


def add_text(slide, l, t, w, h, lines, aligns=None, sizes=None, colors=None, bold=None):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = (aligns or [PP_ALIGN.LEFT]*len(lines))[i]
        r = p.add_run(); r.text = ln
        r.font.size = Pt((sizes or [11]*len(lines))[i])
        r.font.color.rgb = (colors or [INK]*len(lines))[i]
        r.font.bold = (bold or [False]*len(lines))[i]
        r.font.name = 'Arial'
    return tb


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    return s


def grid_slide(prs, title, subtitle, items, cols):
    s = blank(prs)
    add_text(s, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.6),
             [title], [PP_ALIGN.LEFT], [26], [INK], [True])
    add_text(s, Inches(0.52), Inches(0.85), Inches(12.3), Inches(0.35),
             [subtitle], [PP_ALIGN.LEFT], [12], [SUB])
    items = [it for it in items if it[0] and os.path.exists(it[0])]
    n = len(items); rows = (n + cols - 1) // cols
    margin_x, top0 = Inches(0.5), Inches(1.35)
    gap = Inches(0.3)
    cell_w = (Inches(13.33) - 2*margin_x - gap*(cols-1)) / cols
    cap_h = Inches(0.9)
    avail_h = Inches(7.5) - top0 - Inches(0.2)
    cell_h = (avail_h - gap*(rows-1)) / rows
    vid_h = cell_h - cap_h
    for i, (path, ttl, model, inp, mt) in enumerate(items):
        r, c = divmod(i, cols)
        cx = margin_x + c * (cell_w + gap)
        cy = top0 + r * (cell_h + gap)
        pimg, ar = poster(path)
        vw = min(cell_w, Emu(int(vid_h * ar)))
        vx = cx + (cell_w - vw) / 2
        s.shapes.add_movie(path, vx, cy, vw, vid_h, poster_frame_image=pimg, mime_type='video/mp4')
        add_text(s, cx, cy + vid_h + Inches(0.04), cell_w, cap_h,
                 [ttl, model, inp, f'motion: {mt}'],
                 [PP_ALIGN.CENTER]*4, [12, 11, 9, 9],
                 [INK, ACCENT, SUB, SUB], [True, True, False, False])
    return s


def main():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    # title
    s = blank(prs)
    add_text(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.2),
             ['Generated Plant Videos'], [PP_ALIGN.CENTER], [40], [INK], [True])
    add_text(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.2),
             ['Dynamic · Growth · 3D  —  by model & input source',
              'Cosmos-Predict2  ·  Wan 2.2 I2V  ·  TRELLIS  ·  Ours (3DGS + FK physics)'],
             [PP_ALIGN.CENTER]*2, [16, 12], [SUB, ACCENT])
    grid_slide(prs, 'Dynamic — wind sway', 'image→video (Cosmos / Wan) and 3D-render→video (ours)', DYNAMIC, 3)
    grid_slide(prs, 'Growth', 'Cosmos-Predict2  ·  image (realistic seed) → video', GROWTH, 4)
    grid_slide(prs, '3D generation & novel views', 'text→3D→render (TRELLIS)  ·  3D→render (ours)', STATIC3D, 3)
    out = 'outputs/rerun_2026-07/plant_videos.pptx'
    prs.save(out)
    print('saved', out, os.path.getsize(out) // 1024, 'KB')


if __name__ == '__main__':
    main()
